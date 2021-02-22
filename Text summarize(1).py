from webb import webb
from newspaper import Article

from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.stem.snowball import SnowballStemmer
from nltk.stem import WordNetLemmatizer
import nltk

from pptx import slide
from pptx.util import Inches
from pptx import Presentation
from pptx.util import Pt

# from pptx.enum.dml import MSO_THEME_COLOR

k = 0
j = 0

print("1:News\n2:Sports\n3:General");
category = int(input("Select a category(1,2,3 as input): "))  # Getting input for the category

search = list(input("Search: "))  # User Input
print(search)

for i in range(len(search)):
    if search[i] == ' ':
        search[i] = '+'  # Putting '+' sign for the link to work
searchnew = "".join(search)

x = webb.find_all_links('https://www.google.com/search?q=' + searchnew)  # Searching on Google used webb library to
                                                                         # bypass google bot restriction

y = list()
textt = list()

for i in x:
    if category == 1:
        if (k > 10):
            break
        elif (i.startswith('https://')  # Bypassing all the sites we dont want to use in the variable x
              and ('google' not in i)
              and ('history' not in i)
              and ('search' not in i)
              and ('support' not in i)
              and ('twitter' not in i)
              and ('youtube' not in i)
              and ('marketwatch' not in i)):
            if ('theprint' in i) or ('indiatoday' in i) or ('theprint' in i) or ('economictimes' in i) or ('ndtv' in i):
                y.append(i)
                k = k + 1
        else:
            continue

    elif category == 2:
        if (k > 10):
            break
        elif (i.startswith('https://')
              and ('google' not in i)
              and ('history' not in i)
              and ('search' not in i)
              and ('support' not in i)
              and ('twitter' not in i)
              and ('youtube' not in i)
              and ('marketwatch' not in i)):
            if ('firstpost' in i) or ('indiatoday' in i) or ('sportstar' in i) or ('hindustantimes' in i):
                y.append(i)
                k = k + 1
        else:
            continue

    elif category == 3:
        if (k > 5):
            break
        elif (i.startswith('https://')
              and ('google' not in i)
              and ('history' not in i)
              and ('search' not in i)
              and ('support' not in i)
              and ('twitter' not in i)
              and ('youtube' not in i)
              and ('marketwatch' not in i)):
            y.append(i)
            k = k + 1
        else:
            continue

print(y)

textt = ""

for i in y:
    article = Article(i)    # Passing the link to get the Article from the site
    article.download()      # Downloading the Article
    article.parse()
    # print("\n\n\n TITLE \n\n\n")
    # print(article.title)
    # print("\n\n\n ARTICLE \n\n\n")
    # print(article.summary)
    textt = textt + article.text

print()
print()
print("___________________________________________________________________")
print()
print()
# from docx import Document
# from docx.shared import Inches
# import docx NOT python-docx 

# If you get an error uncomment this line and download the necessary libraries
# nltk.download()

#print(textt)

stemmer = SnowballStemmer("english")
stopWords = set(stopwords.words("english"))

words = word_tokenize(textt)
freqTable = dict()
for word in words:
    word = word.lower()
    if word in stopWords:
        continue
    word = stemmer.stem(word)
    if word in freqTable:
        freqTable[word] += 1
    else:
        freqTable[word] = 1
sentences = sent_tokenize(textt)
sentenceValue = dict()
for sentence in sentences:
    for word, freq in freqTable.items():
        if word in sentence.lower():
            if sentence in sentenceValue:
                sentenceValue[sentence] += freq
            else:
                sentenceValue[sentence] = freq

sumValues = 0
for sentence in sentenceValue:
    sumValues += sentenceValue[sentence]

# Average value of a sentence from original text
average = int(sumValues / len(sentenceValue))

summary = ''
temp = ""
sumlist = list()
ct = 1
for sentence in sentences:
    if (sentence in sentenceValue) and (sentenceValue[sentence] > (1.2 * average)):
        summary += " " + sentence

for words in summary.split():
    if ct >= 40 and words.endswith("."):
        temp += (" " + words)
        sumlist.append(temp)
        temp = ""
        ct = 0
    else:
        temp += " " + words
        ct += 1

#print(sumlist)

prs = Presentation()
prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Summary'

tf = body_shape.text_frame
tf.text = sumlist[0]

for i in range(1, len(sumlist)):
    text = "this is a sample text"
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    body_shape = shapes.placeholders[1]

    tf = body_shape.text_frame
    tf.text = sumlist[i]

prs.save('test.pptx')

print(summary)
